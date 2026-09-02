from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Plant Disease Prediction API"
slide.placeholders[1].text = "Presentation for the Flask-based crop disease prediction project"

# Problem statement slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Problem Statement"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "Our smallholder crop farmers face delayed and inaccurate plant disease diagnosis "
    "when they try to identify crop health from photos; current solutions are not enough "
    "because expert advice is scarce, slow, and inaccessible at scale."
)

# Objectives slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Objectives"
body = slide.shapes.placeholders[1].text_frame
body.text = "Main Objective"
p = body.add_paragraph()
p.text = "Develop a Flask API that predicts plant disease from uploaded crop images."
p.level = 1
p.space_after = Pt(12)

p = body.add_paragraph()
p.text = "Specific Objectives"
p.level = 0
p.font.bold = True
p.space_before = Pt(12)

for item in [
    "Create an image upload interface for plant leaf photos.",
    "Use a trained Keras model to classify plant disease from images.",
    "Map model outputs to readable disease labels with labels.txt.",
    "Provide a fast API response for field users and automated systems.",
]:
    p = body.add_paragraph()
    p.text = item
    p.level = 1

# Data analysis slide
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "Data Analysis"
body = slide.shapes.placeholders[1].text_frame
body.text = "Key data analysis steps for the project:"
for item in [
    "Preprocess uploaded images and normalize for model input.",
    "Map predicted indices to human-friendly labels from labels.txt.",
    "Interpret the model's classification results and confidence scores.",
    "Support multi-class disease detection for different crop leaf conditions.",
]:
    p = body.add_paragraph()
    p.text = item
    p.level = 1

# System interface diagram slide
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
slide.shapes.title.text = "System Interface"

left = Inches(0.5)
top = Inches(1.5)
width = Inches(2)
height = Inches(0.8)

shapes = []
labels = ["User / Farmer", "Web Browser", "Flask API / /predict", "Keras Model", "Prediction Result"]
for i, label in enumerate(labels):
    shape = slide.shapes.add_shape(
        1, left + Inches(2.2 * i), top, width, height
    )
    shape.text = label
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(12)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(79, 129, 189)
    shapes.append(shape)

# add arrows connecting shapes
for i in range(len(shapes) - 1):
    x1 = left + Inches(2.2 * i) + width
    y1 = top + height / 2
    x2 = left + Inches(2.2 * (i + 1))
    y2 = y1
    arrow = slide.shapes.add_connector(
        1, x1, y1, x2 - x1, 0
    )
    arrow.line.width = Pt(2)
    arrow.line.color.rgb = RGBColor(0, 0, 0)

prs.save("d:\\Front end\\Project_Presentation.pptx")
print("Saved Project_Presentation.pptx")
